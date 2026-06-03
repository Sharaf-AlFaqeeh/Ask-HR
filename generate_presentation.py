import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_askhr_presentation():
    # 1. Initialize presentation
    prs = Presentation()
    
    # Set to widescreen 16:9 layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 2. Define Theme Colors (Minimalist White & Gray Theme)
    BG_COLOR = RGBColor(255, 255, 255)       # Pure White background
    CARD_COLOR = RGBColor(248, 250, 252)     # Off-white / Very light gray card background (#f8fafc)
    BORDER_COLOR = RGBColor(226, 232, 240)   # Soft light gray border (#e2e8f0)
    TEXT_DARK = RGBColor(15, 23, 42)         # Slate 900 for main titles / headings (#0f172a)
    TEXT_MUTED = RGBColor(51, 65, 85)        # Slate 700 for descriptions (#334155)
    TEXT_LIGHT = RGBColor(100, 116, 139)     # Slate 500 for captions / dates (#64748b)
    
    # Subtle accent colors for statuses/indicators
    GREEN_COLOR = RGBColor(13, 148, 136)     # Soft Teal/Green for AskHR highlights
    RED_COLOR = RGBColor(220, 38, 38)        # Soft Red for Copilot limitations
    GOLD_COLOR = RGBColor(180, 130, 40)      # Muted Gold for Recommendation

    # 3. Helper to set background color
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to add a clean card shape with subtle border
    def add_card(slide, left, top, width, height, fill_color, border_color):
        shape = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE is 1
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
        return shape

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # ==========================================
    # SLIDE 0: TITLE SLIDE
    # ==========================================
    slide_0 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_0, BG_COLOR)

    # Decorative minimalist line on the right side (RTL orientation)
    add_card(slide_0, Inches(12.5), Inches(0.5), Inches(0.1), Inches(6.5), TEXT_DARK, BORDER_COLOR)
    add_card(slide_0, Inches(12.35), Inches(0.5), Inches(0.04), Inches(6.5), TEXT_LIGHT, BORDER_COLOR)

    # Title Text Frame (Aligned Right for Arabic)
    title_box = slide_0.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(10.8), Inches(1.5))
    tf_0 = title_box.text_frame
    tf_0.word_wrap = True
    tf_0.margin_right = Inches(0.2)
    tf_0.margin_left = Inches(0.2)

    p1 = tf_0.paragraphs[0]
    p1.text = "Ask HR Assistant"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_DARK
    p1.alignment = PP_ALIGN.RIGHT

    # Subtitle Text Frame
    subtitle_box = slide_0.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(10.8), Inches(2.2))
    tf_sub = subtitle_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_right = Inches(0.2)
    tf_sub.margin_left = Inches(0.2)

    p2 = tf_sub.add_paragraph()
    p2.text = "مساعد الذكاء الاصطناعي للموارد البشرية لمجموعة هائل سعيد أنعم"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_MUTED
    p2.alignment = PP_ALIGN.RIGHT
    p2.space_before = Pt(8)

    p3 = tf_sub.add_paragraph()
    p3.text = "عرض تقديمي للإدارة حول متطلبات النظام والمقارنة التقنية والمالية مع Copilot"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT_LIGHT
    p3.alignment = PP_ALIGN.RIGHT
    p3.space_before = Pt(12)

    # ==========================================
    # SLIDE 1: SYSTEM REQUIREMENTS (PART 1)
    # ==========================================
    slide_1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_1, BG_COLOR)

    # Top border line
    add_card(slide_1, Inches(1.0), Inches(0.4), Inches(11.33), Inches(0.04), TEXT_DARK, BORDER_COLOR)

    # Slide Title
    title_box_1 = slide_1.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.0))
    tf_1 = title_box_1.text_frame
    tf_1.word_wrap = True
    p_t1 = tf_1.paragraphs[0]
    p_t1.text = "المتطلبات الأساسية للنظام (1) - البنية البرمجية والبيانات"
    p_t1.font.name = "Segoe UI"
    p_t1.font.size = Pt(28)
    p_t1.font.bold = True
    p_t1.font.color.rgb = TEXT_DARK
    p_t1.alignment = PP_ALIGN.RIGHT

    p_sub1 = tf_1.add_paragraph()
    p_sub1.text = "المكونات الأساسية لبناء وتغذية الذكاء الاصطناعي الخاص بالموارد البشرية"
    p_sub1.font.name = "Segoe UI"
    p_sub1.font.size = Pt(14)
    p_sub1.font.color.rgb = TEXT_LIGHT
    p_sub1.alignment = PP_ALIGN.RIGHT
    p_sub1.space_before = Pt(4)

    # Content Columns (RTL: Column 1 on Right, Column 2 on Left)
    
    # Column 1 (Right): Software Infrastructure
    add_card(slide_1, Inches(7.0), Inches(1.8), Inches(5.33), Inches(4.8), CARD_COLOR, BORDER_COLOR)
    card_tf1 = slide_1.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.33), Inches(4.8)).text_frame
    card_tf1.word_wrap = True
    card_tf1.margin_left = Inches(0.3)
    card_tf1.margin_right = Inches(0.3)
    card_tf1.margin_top = Inches(0.3)
    card_tf1.margin_bottom = Inches(0.3)
    
    p_c1_t = card_tf1.paragraphs[0]
    p_c1_t.text = "⚙️ الأساس البرمجي والتشغيل"
    p_c1_t.font.name = "Segoe UI"
    p_c1_t.font.size = Pt(20)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = TEXT_DARK
    p_c1_t.alignment = PP_ALIGN.RIGHT
    
    points1 = [
        ("لغة بايثون (Python)", "الأساس البرمجي لربط الخدمات ومعالجة البيانات وبناء منطق العمل والنظام الرئيسي عبر FastAPI."),
        ("نموذج لغوي مدرب مسبقاً", "استخدام نموذج لغوي مفتوح المصدر (مثل Qwen GGUF) مهيأ للاستدعاء المحلي والاستدلال المباشر."),
        ("إضافات ومكتبات تشغيل النموذج", "مكتبات معالجة اللغات الطبيعية والاستخلاص المحلي للكيانات (llama-cpp-python و FastEmbed).")
    ]
    for title, desc in points1:
        p_title = card_tf1.add_paragraph()
        p_title.text = f"• {title}"
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        p_title.alignment = PP_ALIGN.RIGHT
        p_title.space_before = Pt(12)
        
        p_desc = card_tf1.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Segoe UI"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.alignment = PP_ALIGN.RIGHT
        p_desc.space_before = Pt(3)

    # Column 2 (Left): Data & Vector Integration (RAG)
    add_card(slide_1, Inches(1.0), Inches(1.8), Inches(5.33), Inches(4.8), CARD_COLOR, BORDER_COLOR)
    card_tf2 = slide_1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.33), Inches(4.8)).text_frame
    card_tf2.word_wrap = True
    card_tf2.margin_left = Inches(0.3)
    card_tf2.margin_right = Inches(0.3)
    card_tf2.margin_top = Inches(0.3)
    card_tf2.margin_bottom = Inches(0.3)
    
    p_c2_t = card_tf2.paragraphs[0]
    p_c2_t.text = "📂 تغذية البيانات والمعرفة (RAG)"
    p_c2_t.font.name = "Segoe UI"
    p_c2_t.font.size = Pt(20)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = TEXT_DARK
    p_c2_t.alignment = PP_ALIGN.RIGHT
    
    points2 = [
        ("ملفات سياسات المجموعة (PDF/Docs)", "اللوائح الداخلية، السياسات الإدارية، وإجراءات العمل المعتمدة في مجموعة هائل سعيد أنعم."),
        ("بيانات قسم الموارد البشرية (FAQ)", "تجميع وتجهيز الأسئلة الأكثر شيوعاً وإجاباتها الرسمية لتوفير مرجعية سريعة."),
        ("قاعدة بيانات المتجهات (Vector DB)", "استخدام Qdrant محلياً لتخزين متجهات النصوص واسترجاع المعلومات المناسبة للرد بدقة عالية (RAG).")
    ]
    for title, desc in points2:
        p_title = card_tf2.add_paragraph()
        p_title.text = f"• {title}"
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        p_title.alignment = PP_ALIGN.RIGHT
        p_title.space_before = Pt(12)
        
        p_desc = card_tf2.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Segoe UI"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.alignment = PP_ALIGN.RIGHT
        p_desc.space_before = Pt(3)

    # ==========================================
    # SLIDE 2: SYSTEM REQUIREMENTS (PART 2)
    # ==========================================
    slide_2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_2, BG_COLOR)

    # Top border line
    add_card(slide_2, Inches(1.0), Inches(0.4), Inches(11.33), Inches(0.04), TEXT_DARK, BORDER_COLOR)

    # Slide Title
    title_box_2 = slide_2.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.0))
    tf_2 = title_box_2.text_frame
    tf_2.word_wrap = True
    p_t2 = tf_2.paragraphs[0]
    p_t2.text = "المتطلبات الأساسية للنظام (2) - بيئات العمل ومتطلبات التشغيل"
    p_t2.font.name = "Segoe UI"
    p_t2.font.size = Pt(28)
    p_t2.font.bold = True
    p_t2.font.color.rgb = TEXT_DARK
    p_t2.alignment = PP_ALIGN.RIGHT

    p_sub2 = tf_2.add_paragraph()
    p_sub2.text = "من مرحلة التطوير والاختبار المحلي إلى مرحلة الاستضافة والإنتاج الفعلي"
    p_sub2.font.name = "Segoe UI"
    p_sub2.font.size = Pt(14)
    p_sub2.font.color.rgb = TEXT_LIGHT
    p_sub2.alignment = PP_ALIGN.RIGHT
    p_sub2.space_before = Pt(4)

    # Content Columns (RTL: Column 1 on Right, Column 2 on Left)

    # Column 1 (Right): Development Environment (Laptop)
    add_card(slide_2, Inches(7.0), Inches(1.8), Inches(5.33), Inches(4.8), CARD_COLOR, BORDER_COLOR)
    dev_tf = slide_2.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.33), Inches(4.8)).text_frame
    dev_tf.word_wrap = True
    dev_tf.margin_left = Inches(0.3)
    dev_tf.margin_right = Inches(0.3)
    dev_tf.margin_top = Inches(0.3)
    dev_tf.margin_bottom = Inches(0.3)

    p_dev_t = dev_tf.paragraphs[0]
    p_dev_t.text = "💻 بيئة التطوير والتجربة الأولوية"
    p_dev_t.font.name = "Segoe UI"
    p_dev_t.font.size = Pt(20)
    p_dev_t.font.bold = True
    p_dev_t.font.color.rgb = TEXT_DARK
    p_dev_t.alignment = PP_ALIGN.RIGHT

    dev_points = [
        ("تطوير محلي على لابتوب شخصي", "الهدف هو السماح للمطورين ببناء الميزات واختبارها محلياً دون الحاجة لاتصال سيرفر مستمر."),
        ("استخدام نماذج مكممة (Quantized)", "تشغيل إصدارات GGUF خفيفة من النماذج اللغوية لتناسب موارد اللابتوب الشخصي وتسهيل عملية التطوير السريع."),
        ("أدوات الفحص والتحقق الآلي", "تكامل شامل مع pytest لاختبار دقة استخراج البيانات وصحة تدفق الحوارات محلياً.")
    ]
    for title, desc in dev_points:
        p_title = dev_tf.add_paragraph()
        p_title.text = f"• {title}"
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        p_title.alignment = PP_ALIGN.RIGHT
        p_title.space_before = Pt(12)
        
        p_desc = dev_tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Segoe UI"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.alignment = PP_ALIGN.RIGHT
        p_desc.space_before = Pt(3)

    # Column 2 (Left): Production Server Environment (On-Premises)
    add_card(slide_2, Inches(1.0), Inches(1.8), Inches(5.33), Inches(4.8), CARD_COLOR, BORDER_COLOR)
    prod_tf = slide_2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.33), Inches(4.8)).text_frame
    prod_tf.word_wrap = True
    prod_tf.margin_left = Inches(0.3)
    prod_tf.margin_right = Inches(0.3)
    prod_tf.margin_top = Inches(0.3)
    prod_tf.margin_bottom = Inches(0.3)

    p_prod_t = prod_tf.paragraphs[0]
    p_prod_t.text = "🖥️ خادم التشغيل المحلي (Production)"
    p_prod_t.font.name = "Segoe UI"
    p_prod_t.font.size = Pt(20)
    p_prod_t.font.bold = True
    p_prod_t.font.color.rgb = TEXT_DARK
    p_prod_t.alignment = PP_ALIGN.RIGHT

    prod_points = [
        ("معالج كرت شاشة (GPU) مخصص", "كرت شاشة بذاكرة VRAM تتراوح من 8 إلى 16 جيجابايت (مثل RTX 3090/4080 أو كروت فئة سيرفر A10G/A30)."),
        ("تسريع الاستجابة والطلبات المتزامنة", "الـ VRAM الكبيرة تسمح للنموذج بمعالجة سياقات أطول (كبيانات سياسات HR) وتحسين المعالجة المتزامنة لآلاف الموظفين."),
        ("استقلالية تامة وبيانات معزولة", "استضافة الخادم داخلياً يضمن الامتثال لسياسات أمن المعلومات للمجموعة وعدم خروج البيانات الحساسة.")
    ]
    for title, desc in prod_points:
        p_title = prod_tf.add_paragraph()
        p_title.text = f"• {title}"
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        p_title.alignment = PP_ALIGN.RIGHT
        p_title.space_before = Pt(12)
        
        p_desc = prod_tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Segoe UI"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.alignment = PP_ALIGN.RIGHT
        p_desc.space_before = Pt(3)


    # ==========================================
    # SLIDE 3: COMPARISON WITH COPILOT
    # ==========================================
    slide_3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_3, BG_COLOR)

    # Top border line
    add_card(slide_3, Inches(1.0), Inches(0.4), Inches(11.33), Inches(0.04), TEXT_DARK, BORDER_COLOR)

    # Slide Title
    title_box_3 = slide_3.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.0))
    tf_3 = title_box_3.text_frame
    tf_3.word_wrap = True
    p_t3 = tf_3.paragraphs[0]
    p_t3.text = "المقارنة مع Copilot - الجدوى الاقتصادية والتقنية والسيادة الأمنية"
    p_t3.font.name = "Segoe UI"
    p_t3.font.size = Pt(28)
    p_t3.font.bold = True
    p_t3.font.color.rgb = TEXT_DARK
    p_t3.alignment = PP_ALIGN.RIGHT

    p_sub3 = tf_3.add_paragraph()
    p_sub3.text = "مقارنة شاملة بين ترخيص Copilot السحابي والنظام الداخلي المطور محلياً لآلاف الموظفين"
    p_sub3.font.name = "Segoe UI"
    p_sub3.font.size = Pt(14)
    p_sub3.font.color.rgb = TEXT_LIGHT
    p_sub3.alignment = PP_ALIGN.RIGHT
    p_sub3.space_before = Pt(4)

    # Three-Column Comparison Layout (RTL: Column 1 on Right, Column 2 in Middle, Column 3 on Left)
    
    # Column 1 (Right): Custom AskHR (Local RAG)
    add_card(slide_3, Inches(9.0), Inches(1.8), Inches(3.33), Inches(4.8), CARD_COLOR, BORDER_COLOR)
    # Highlight top border of custom solution card
    add_card(slide_3, Inches(9.0), Inches(1.8), Inches(3.33), Inches(0.08), GREEN_COLOR, GREEN_COLOR)
    col1_tf = slide_3.shapes.add_textbox(Inches(9.0), Inches(1.9), Inches(3.33), Inches(4.7)).text_frame
    col1_tf.word_wrap = True
    col1_tf.margin_left = Inches(0.2)
    col1_tf.margin_right = Inches(0.2)
    col1_tf.margin_top = Inches(0.2)
    col1_tf.margin_bottom = Inches(0.2)

    p_col1_t = col1_tf.paragraphs[0]
    p_col1_t.text = "✅ نظام AskHR الخاص"
    p_col1_t.font.name = "Segoe UI"
    p_col1_t.font.size = Pt(18)
    p_col1_t.font.bold = True
    p_col1_t.font.color.rgb = GREEN_COLOR
    p_col1_t.alignment = PP_ALIGN.RIGHT

    askhr_points = [
        ("صفر تكلفة تراخيص مستمرة", "تطوير أولي لمرة واحدة (CapEx) يعقبه تكاليف خادم محلي ثابتة وبسيطة (OpEx). تكلفة الموظف الإضافي أو الاستخدام الإضافي = صفر."),
        ("تحكم وتخصيص مرن 100%", "بني بمعمارية Clean Ports & Adapters للربط العميق والكامل مع SAP SF وقواعد البيانات بدون قيود أطر العمل الجاهزة."),
        ("سيادة أمنية مطلقة (On-Prem)", "تشغيل نموذج محلي 100% يحمي سرية معلومات الرواتب والإجازات ويبقى داخل شبكة المجموعة الآمنة.")
    ]
    for title, desc in askhr_points:
        p_title = col1_tf.add_paragraph()
        p_title.text = f"• {title}"
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        p_title.alignment = PP_ALIGN.RIGHT
        p_title.space_before = Pt(10)
        
        p_desc = col1_tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Segoe UI"
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.alignment = PP_ALIGN.RIGHT
        p_desc.space_before = Pt(3)

    # Column 2 (Middle): Microsoft Copilot
    add_card(slide_3, Inches(5.33), Inches(1.8), Inches(3.33), Inches(4.8), CARD_COLOR, BORDER_COLOR)
    # Highlight top border of copilot card
    add_card(slide_3, Inches(5.33), Inches(1.8), Inches(3.33), Inches(0.08), RED_COLOR, RED_COLOR)
    col2_tf = slide_3.shapes.add_textbox(Inches(5.33), Inches(1.9), Inches(3.33), Inches(4.7)).text_frame
    col2_tf.word_wrap = True
    col2_tf.margin_left = Inches(0.2)
    col2_tf.margin_right = Inches(0.2)
    col2_tf.margin_top = Inches(0.2)
    col2_tf.margin_bottom = Inches(0.2)

    p_col2_t = col2_tf.paragraphs[0]
    p_col2_t.text = "❌ خيار Microsoft Copilot"
    p_col2_t.font.name = "Segoe UI"
    p_col2_t.font.size = Pt(18)
    p_col2_t.font.bold = True
    p_col2_t.font.color.rgb = RED_COLOR
    p_col2_t.alignment = PP_ALIGN.RIGHT

    copilot_points = [
        ("تكلفة مرتفعة ومستمرة", "تراخيص شهرية ($30/مستخدم لـ M365) أو رسوم رسائل إضافية في Copilot Studio ($200 لـ 2000 رسالة). تزيد مع زيادة عدد الموظفين."),
        ("محدود التخصيص والربط", "صعوبة بالغة في تخصيص تدفقات الحوار المعقدة (Slot-filling) والربط المباشر مع أنظمة مخصصة مثل SAP SF بدون رسوم إضافية."),
        ("سحابي وخطر أمني", "يتم إرسال بيانات الموظفين ورواتبهم الحساسة إلى خادم سحابي خارجي، مما يتعارض مع سياسات حوكمة البيانات.")
    ]
    for title, desc in copilot_points:
        p_title = col2_tf.add_paragraph()
        p_title.text = f"• {title}"
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        p_title.alignment = PP_ALIGN.RIGHT
        p_title.space_before = Pt(10)
        
        p_desc = col2_tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Segoe UI"
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.alignment = PP_ALIGN.RIGHT
        p_desc.space_before = Pt(3)

    # Column 3 (Left): Managerial Conclusion & Recommendation (ترجيح الخيار)
    add_card(slide_3, Inches(1.0), Inches(1.8), Inches(4.0), Inches(4.8), CARD_COLOR, BORDER_COLOR)
    # Highlight top border of recommendation card
    add_card(slide_3, Inches(1.0), Inches(1.8), Inches(4.0), Inches(0.08), GOLD_COLOR, GOLD_COLOR)
    col3_tf = slide_3.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(4.0), Inches(4.7)).text_frame
    col3_tf.word_wrap = True
    col3_tf.margin_left = Inches(0.25)
    col3_tf.margin_right = Inches(0.25)
    col3_tf.margin_top = Inches(0.25)
    col3_tf.margin_bottom = Inches(0.25)

    p_col3_t = col3_tf.paragraphs[0]
    p_col3_t.text = "🎯 التوصية الإدارية والمالية"
    p_col3_t.font.name = "Segoe UI"
    p_col3_t.font.size = Pt(18)
    p_col3_t.font.bold = True
    p_col3_t.font.color.rgb = GOLD_COLOR
    p_col3_t.alignment = PP_ALIGN.RIGHT

    rec_title = col3_tf.add_paragraph()
    rec_title.text = "ترجيح خيار: النظام الخاص المطور محلياً"
    rec_title.font.name = "Segoe UI"
    rec_title.font.size = Pt(14)
    rec_title.font.bold = True
    rec_title.font.color.rgb = TEXT_DARK
    rec_title.alignment = PP_ALIGN.RIGHT
    rec_title.space_before = Pt(12)

    rec_desc = col3_tf.add_paragraph()
    rec_desc.text = (
        "عند حساب التكلفة التشغيلية لآلاف الموظفين، نجد أن العائد الاستثماري (ROI) للنظام الخاص يتحقق خلال أقل من 3 أشهر مقارنة بكوبايلوت. "
        "الاستثمار في سيرفر محلي (CapEx) بقيمة 3000-5000 دولار لمرة واحدة هو حل دائم ومستقر ويلغي تماماً الاشتراك الشهري المتكرر البالغ آلاف الدولارات. "
        "\n\nعلاوة على ذلك، حماية بيانات الموظفين محلياً وإمكانية تخصيص الحوار ليتلاءم مع العمليات الدقيقة في SAP SF تجعل النظام الخاص هو الخيار الاستراتيجي الأنسب للمجموعة."
    )
    rec_desc.font.name = "Segoe UI"
    rec_desc.font.size = Pt(11)
    rec_desc.font.color.rgb = TEXT_MUTED
    rec_desc.alignment = PP_ALIGN.RIGHT
    rec_desc.space_before = Pt(8)

    # 4. Save presentation
    output_filename = "AskHR_Presentation.pptx"
    output_path = os.path.join("p:\\____AI____\\HSAGroup\\AskHRPro", output_filename)
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_askhr_presentation()
